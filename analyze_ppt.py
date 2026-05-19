from pptx import Presentation

try:
    prs = Presentation("presentation.pptx")
    print(f"Loaded 'presentation.pptx' successfully. Slides: {len(prs.slides)}")
    
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        if slide.shapes.title:
            print(f"Title: {slide.shapes.title.text}")
        else:
            print("Title: [No Title]")
        
        print("Shapes:")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(f" - {shape.shape_type}: {shape.text[:50]}...")
            else:
                print(f" - {shape.shape_type}")

except Exception as e:
    print(f"Error loading presentation: {e}")
