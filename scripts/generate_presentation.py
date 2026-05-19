from pptx import Presentation
from pptx.util import Inches, Pt
import os
import datetime

SOURCE_PPT = "presentation.pptx"
OUTPUT_PPT = "Smart_Pest_Management_System.pptx"
SCREENSHOT_DIR = "screenshots"

REPLACEMENTS = {
    "FarmEco Assist AI": "Smart Pest Management System",
    "January 12, 2026": datetime.datetime.now().strftime("%B %d, %Y"),
    "CSD334": "CSD400", # Example Code
    "VISMAYA JOSEPH": "Project Team",
    "MBI23CS136": "",
    "Sameena K M": "Project Guide",
    "Assistant Professor": "",
    "CSE Dept.": "CSE Department",
}

def replace_text_in_shape(shape, replacements):
    if not shape.has_text_frame:
        return
    
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            original_text = run.text
            for old, new in replacements.items():
                if old in original_text:
                    run.text = original_text.replace(old, new)

def generate():
    if not os.path.exists(SOURCE_PPT):
        print(f"Error: {SOURCE_PPT} not found.")
        return

    print(f"Loading {SOURCE_PPT}...")
    prs = Presentation(SOURCE_PPT)

    # 1. Global Text Replacement
    print("Performing text replacements...")
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_text_in_shape(shape, REPLACEMENTS)

    # 2. Update specific slides
    # Slide 11, 12, 13 are for screenshots (indices 10, 11, 12)
    screenshot_map = {
        10: "landing_page.png",
        11: "farmer_dashboard.png",
        12: "admin_dashboard.png" # Or expert
    }

    for slide_idx, img_name in screenshot_map.items():
        if slide_idx < len(prs.slides):
            slide = prs.slides[slide_idx]
            img_path = os.path.join(SCREENSHOT_DIR, img_name)
            
            if os.path.exists(img_path):
                print(f"Updating Slide {slide_idx+1} with {img_name}...")
                
                # Find the main picture placeholder. 
                # Strategy: Find any existing picture, remove it, and place new one in similar position
                # Or if there is a placeholder, use it.
                
                # Let's try to find a picture shape
                pic_shape = None
                for shape in slide.shapes:
                    if shape.shape_type == 13: # PICTURE
                        pic_shape = shape
                        break
                
                if pic_shape:
                    # Capture geometry
                    left, top, width, height = pic_shape.left, pic_shape.top, pic_shape.width, pic_shape.height
                    # Remove old picture (trickier in python-pptx, usually we remove from shape tree)
                    sp = pic_shape._element
                    sp.getparent().remove(sp)
                    
                    # Add new picture
                    slide.shapes.add_picture(img_path, left, top, width, height)
                else:
                    # Just add centered
                    slide.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(8))
            else:
                print(f"Warning: {img_name} not found.")

    # Save
    prs.save(OUTPUT_PPT)
    print(f"Presentation saved to {OUTPUT_PPT}")

if __name__ == "__main__":
    generate()
