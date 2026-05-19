from pptx import Presentation
from pptx.util import Inches, Pt
import os
import datetime

SOURCE_PPT = "presentation.pptx"
OUTPUT_PPT = "Corrected_Pest_System_Presentation.pptx"
SCREENSHOT_DIR = "screenshots"

# Mapping specific slides (indices) to new content (Title, Bullets)
# Index is 0-based
SLIDE_UPDATES = {
    2: { # Slide 3: INTRODUCTION
        "title": "INTRODUCTION",
        "content": [
            "Smart Pest Management System: An AI-driven solution for modern farming.",
            "Bridging the gap between technology and agriculture.",
            "Empowers farmers with instant plant disease detection.",
            "Provides organic and inorganic treatment recommendations.",
            "Connects farmers with experts and local pesticide shops."
        ]
    },
    3: { # Slide 4: PROBLEM STATEMENT
        "title": "PROBLEM STATEMENT",
        "content": [
            "Delayed disease identification leads to significant crop loss (up to 40%).",
            "Lack of accessible expert advice in remote areas.",
            "Overuse of harmful chemical pesticides due to lack of knowledge.",
            "Difficulty in finding authentic pesticide shops.",
            "Existing solutions often lack regional context or offline capabilities."
        ]
    },
    4: { # Slide 5: SYSTEM ARCHITECTURE (Modules)
        "title": "SYSTEM MODULES",
        "content": [
            "1. Farmer Module: Disease Detection, Weather Updates, Shop Locator, Chatbot.",
            "2. Admin Module: User Management, Content Management (Shops/Medicines).",
            "3. Expert Module: Consultation Dashboard for answering farmer queries.",
            "Integrated via a central MongoDB database and Flask API."
        ]
    },
    5: { # Slide 6: Tech Specs
        "title": "TECHNICAL SPECIFICATIONS",
        "content": [
            "Frontend: React.js v19, Vite, Tailwind CSS (Modern UI)",
            "Backend: Python Flask, PyTorch (MobileNetV2)",
            "Database: MongoDB (NoSQL) with PyMongo",
            "AI Model: MobileNetV2 (Transfer Learning on PlantVillage)",
            "Deployment: Localhost / Clean Architecture"
        ]
    },
    6: { # Slide 7: DATA FLOW
        "title": "DATA FLOW",
        "content": [
            "1. User Uploads Image (Frontend)",
            "2. Image sent to Flask Backend (API)",
            "3. Pre-processed (Resize/Normalize) & passed to PyTorch Model",
            "4. Classification Result + Confidence Score returned",
            "5. System fetches Treatment Data from MongoDB",
            "6. Full Diagnosis displayed to User"
        ]
    },
    7: { # Slide 8: SYSTEM WORKFLOW
        "title": "SYSTEM WORKFLOW",
        "content": [
            "Start -> Login/Register",
            "Upload Leaf Image -> Analyze",
            "View Result (Disease Name, Confidence, Treatments)",
            "Optional: Chat with Bot or Request Expert Consultation",
            "Optional: Locate nearby Shops or View Subsidies",
            "End"
        ]
    },
    8: { # Slide 9: Code Structure
        "title": "IMPLEMENTATION – Code Structure",
        "content": [
            "/src - React Function Components (Hooks, Context)",
            "/backend - Flask Application Entry Point",
            "/backend/data - Model Weights & Datasets",
            "/components - Reusable UI Elements"
        ]
    },
    9: { # Slide 10: Software Architecture
        "title": "SOFTWARE ARCHITECTURE",
        "content": [
            "Client-Server Architecture (REST API)",
            "Frontend sends image/data via HTTP Requests",
            "Flask Backend processes images and queries MongoDB",
            "AI Inference runs locally using PyTorch",
            "JWT/Session-based Authentication"
        ]
    },
    13: { # Slide 14: FEATURES
        "title": "KEY FEATURES",
        "content": [
            "AI Disease Detection: High accuracy using CNN.",
            "Treatment Filtering: Color-coded (Organic=Green, Chemical=Red).",
            "Offline Chatbot: Resource for farmers without internet.",
            "Shop Locator: Google Maps-like directory.",
            "Community Forum: Peer-to-peer farmer support.",
            "Government Subsidies Portal."
        ]
    },
    14: { # Slide 15: Progress
        "title": "PROJECT STATUS",
        "content": [
            "Current Status: 100% Core Features Completed",
            "Implemented Features:",
            "- Disease Detection & Medicine Recommendation",
            "- Admin & Expert Dashboards",
            "- User Management & Shop Locator",
            "- Offline Chatbot & Forum"
        ]
    },
    15: { # Slide 16: FUTURE ROADMAP
        "title": "FUTURE ROADMAP",
        "content": [
            "Mobile App Development (React Native).",
            "Drone Integration for large-scale field scanning.",
            "Multilingual Support (Regional Languages).",
            "Real-time Weather Integration API.",
            "Marketplace for buying/selling crops directly."
        ]
    },
    16: { # Slide 17: CONCLUSION
        "title": "CONCLUSION",
        "content": [
            "The Smart Pest Management System effectively addresses the challenges of modern farming.",
            "It provides timely, accurate, and accessible solutions.",
            "Promotes sustainable agricultural practices.",
            "Bridged the digital divide for farmers.",
            "Ready for real-world deployment and testing."
        ]
    }
}

REPLACEMENTS = {
    "FarmEco Assist AI": "Smart Pest Management System",
    "January 12, 2026": datetime.datetime.now().strftime("%B %d, %Y"),
    "CSD334": "CSD400",
    "Streamlit": "React & Flask",
    "VISMAYA JOSEPH": "Project Team",
    "MBI23CS136": "",
    "Sameena K M": "Project Guide",
    "Assistant Professor": "",
    "CSE Dept.": "CSE Department",
}

def clear_and_set_content(slide, title_text, bullets):
    # Set Title
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    
    # helper to find the main text body. Usually the second shape or a placeholder.
    body_shape = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1: # Often the body
            body_shape = shape
            break
            
    # If no placeholder, look for a large text box
    if not body_shape:
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                body_shape = shape
                break
    
    if body_shape and body_shape.has_text_frame:
        tf = body_shape.text_frame
        tf.clear() # Wipe old content
        
        for point in bullets:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            # Optional: Add sub-bullets if needed
            
def replace_text_general(shape, replacements):
    if not shape.has_text_frame: return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            original = run.text
            for old, new in replacements.items():
                if old in original:
                    run.text = original.replace(old, new)

def generate():
    if not os.path.exists(SOURCE_PPT):
        print(f"Error: {SOURCE_PPT} found.")
        return

    print("Loading presentation...")
    prs = Presentation(SOURCE_PPT)

    # 1. Slide-Specific Overwrites (Logic First)
    for index, data in SLIDE_UPDATES.items():
        if index < len(prs.slides):
            print(f"Overwriting Slide {index+1}...")
            clear_and_set_content(prs.slides[index], data['title'], data['content'])

    # 2. General Replacements (for remaining slides, e.g. Title, Intro)
    # i=0 is title slide, mostly handled by REPLACEMENTS
    for i, slide in enumerate(prs.slides):
        # We process ALL slides for variable replacements, even those we overwrote, just in case
        for shape in slide.shapes:
            replace_text_general(shape, REPLACEMENTS)

    # 3. Insert Screenshots (Slides 11, 12, 13 -> Indices 10, 11, 12)
    # Note: If these slides were overwritten by SLIDE_UPDATES, we need to be careful.
    # But SLIDE_UPDATES keys above are 2,3,4,5,6,7,8,9,13,14,15,16.
    # Indices 10, 11, 12 (Slides 11, 12, 13) are NOT in SLIDE_UPDATES so they are safe to just inject images.
    
    current_screenshot_idx = 0
    screenshots = [
        "landing_page.png",
        "farmer_dashboard.png", 
        "admin_dashboard.png",
        "expert_dashboard.png" 
    ]
    
    # Target indices for screenshots
    target_slides = [10, 11, 12] # Slides 11, 12, 13
    
    for i, slide_idx in enumerate(target_slides):
        if slide_idx < len(prs.slides) and i < len(screenshots):
            img_name = screenshots[i]
            img_path = os.path.join(SCREENSHOT_DIR, img_name)
            
            if os.path.exists(img_path):
                print(f"Injecting {img_name} into Slide {slide_idx+1}")
                slide = prs.slides[slide_idx]
                
                # Try to remove existing pictures
                for shape in list(slide.shapes):
                    if shape.shape_type == 13: # PICTURE
                        sp = shape._element
                        sp.getparent().remove(sp)
                
                # Add new picture (Centered-ish)
                slide.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(8))

    prs.save(OUTPUT_PPT)
    print(f"Saved corrected presentation to {OUTPUT_PPT}")

if __name__ == "__main__":
    generate()
